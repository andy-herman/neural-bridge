"""Extract INFO 310A lecture .pptx files into the vault corpus.

Source:  ~/Desktop/Andy Herman/INFO 310A SP 2026 (Updated)/Lecture_<N>_*.pptx
Target:  ~/Documents/Luna Master/Neural Bridge/Corpus/INFO 310A/lectures/<NN>-<slug>.md

Each output file has TWO sections:
  (A) Full extract — overwritten on every re-run (source of truth from .pptx)
  (B) Research dossier — preserved across re-runs (populated by professor agent)

Usage:
    ~/Development/neural-bridge/.venv/bin/python scripts/corpus/extract_lectures.py
    # or for a single lecture:
    ~/Development/neural-bridge/.venv/bin/python scripts/corpus/extract_lectures.py --only 7

Idempotent: re-running with updated source .pptx files overwrites Section A,
preserves Section B verbatim.

Skip range: 16-20 (per Andy's directive — those are being rewritten separately).
"""

from __future__ import annotations

import argparse
import re
import sys
from datetime import date
from pathlib import Path

from pptx import Presentation

SOURCE_DIR = Path.home() / "Desktop" / "Andy Herman" / "INFO 310A SP 2026 (Updated)"
VAULT_LECTURES_DIR = (
    Path.home() / "Documents" / "Luna Master" / "Neural Bridge" / "Corpus" / "INFO 310A" / "lectures"
)

# Topic slug + display title per lecture number (from Canvas syllabus).
# Lectures 16-20 intentionally omitted — Andy is rewriting those.
LECTURE_META = {
    1:  ("welcome-and-intro",                       "Welcome and Introduction to Cybersecurity",  1, "2026-03-31"),
    2:  ("ai-risk-management",                      "AI Risk Management",                          1, "2026-04-02"),
    3:  ("networking-osi-mitm",                     "Intro to Networking, the OSI Model, and MITM", 2, "2026-04-07"),
    4:  ("networking-devices-subnetting",           "Intro to Networking Devices and Subnetting", 2, "2026-04-09"),
    5:  ("cryptography-intro",                      "Intro to Cryptography",                       3, "2026-04-14"),
    6:  ("cryptography-public-key",                 "Public Key Cryptography and PKI",             3, "2026-04-16"),
    7:  ("application-layer-architecture",          "The Application Layer, Architecture, Logging", 4, "2026-04-21"),
    8:  ("application-frontend-backend-3rd-parties","Application Frontend, Backend, and 3rd Parties", 4, "2026-04-23"),
    9:  ("sessions-and-authentication",             "Sessions and Authentication",                 5, "2026-04-28"),
    10: ("authorization-and-permissions",           "Authorization and Permissions",               6, "2026-05-05"),
    11: ("hackers-social-engineering-malware",      "Hackers, Social Engineering, Malware, and Offensive Security", 6, "2026-05-07"),
    12: ("threat-modeling-stride-dread",            "Threat Modeling, STRIDE, and DREAD",          7, "2026-05-12"),
    13: ("owasp-injection",                         "OWASP Top 10 / Injection",                    7, "2026-05-14"),
    14: ("xss-sast-dast",                           "XSS and Testing Web Services with SAST and DAST", 8, "2026-05-19"),
    15: ("manual-testing-and-hardening",            "Manual Testing and Hardening",                8, "2026-05-21"),
}

# Filename → lecture number. We accept any of the actual filenames in the source dir.
FILENAME_RE = re.compile(r"^Lecture_(\d+)_.*\.pptx$")

SECTION_B_HEADER = "## (B) Research dossier"
SECTION_A_HEADER = "## (A) Full extract"


def find_source_pptx(lecture_n: int) -> Path | None:
    """Find the .pptx for a given lecture number in SOURCE_DIR."""
    candidates = []
    for p in SOURCE_DIR.glob(f"Lecture_{lecture_n}_*.pptx"):
        # Filter out duplicates ending in _Accessible.pptx vs the bare version —
        # prefer _Accessible if both present, but the cleaned-up folder should
        # have only one per lecture now.
        candidates.append(p)
    if not candidates:
        return None
    if len(candidates) == 1:
        return candidates[0]
    # If both versions still exist, prefer _Accessible (the canonical one).
    accessible = [c for c in candidates if "_Accessible" in c.name]
    if accessible:
        return accessible[0]
    return candidates[0]


def extract_slide(slide_idx: int, slide) -> str:
    """Render one slide as markdown."""
    title = ""
    body_parts: list[str] = []

    # Title is usually the first placeholder with placeholder_format.idx == 0
    # but defensively scan all shapes.
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        title = slide.shapes.title.text_frame.text.strip()

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # Skip the title shape (already captured)
        if slide.shapes.title is not None and shape == slide.shapes.title:
            continue
        text = shape.text_frame.text.strip()
        if text:
            body_parts.append(text)

    # Speaker notes
    notes = ""
    if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
        notes = slide.notes_slide.notes_text_frame.text.strip()

    out = [f"### Slide {slide_idx}" + (f": {title}" if title else "")]
    if body_parts:
        out.append("**Body:**")
        out.append("")
        for part in body_parts:
            # Render multi-line text frames as block quote so structure is preserved
            for line in part.splitlines():
                line = line.rstrip()
                if line:
                    out.append(f"> {line}")
                else:
                    out.append(">")
        out.append("")
    if notes:
        out.append("**Speaker notes:**")
        out.append("")
        for line in notes.splitlines():
            line = line.rstrip()
            if line:
                out.append(line)
        out.append("")
    return "\n".join(out)


def extract_pptx(path: Path) -> tuple[int, str]:
    """Extract a .pptx into the section-A markdown body. Returns (slide_count, markdown)."""
    prs = Presentation(str(path))
    blocks = []
    for i, slide in enumerate(prs.slides, start=1):
        blocks.append(extract_slide(i, slide))
    return len(prs.slides), "\n\n".join(blocks)


def render_section_a(lecture_n: int, source_path: Path, slide_count: int, body: str) -> str:
    today = date.today().isoformat()
    return (
        f"{SECTION_A_HEADER}\n\n"
        f"> Extracted from `{source_path.name}` on {today}. {slide_count} slides. "
        f"Re-run `scripts/corpus/extract_lectures.py` to refresh from the source .pptx.\n\n"
        f"{body}\n"
    )


def render_section_b_placeholder() -> str:
    return (
        f"{SECTION_B_HEADER}\n\n"
        f"> Populated by the `teaching-prep` (professor) agent. Re-extraction does NOT overwrite this section.\n"
        f"> Last updated: _(not yet populated — run a professor research pass to fill)_\n\n"
        f"### Accuracy check\n\n"
        f"_What in this lecture is correct, what's debatable, what's drifted from current consensus._\n\n"
        f"### Modern security trends to incorporate\n\n"
        f"_What's hot in this topic right now (2026) that the slides should mention._\n\n"
        f"### Easier ways to land each concept\n\n"
        f"_Specific concepts the professor flags as hard for students; alternative framings._\n\n"
        f"### Real-world story suggestions\n\n"
        f"_Concrete recent incidents Andy could weave in to make the concept stick._\n\n"
        f"### Terminal-novice friction points\n\n"
        f"_Where this lecture's content assumes terminal/Docker/git fluency students don't have._\n"
    )


def render_frontmatter(lecture_n: int, slug: str, title: str, week: int, lecture_date: str,
                       source_path: Path, slide_count: int) -> str:
    return (
        "---\n"
        f"type: lecture-corpus\n"
        f"lecture_number: {lecture_n}\n"
        f"slug: {slug}\n"
        f"title: {title!r}\n"
        f"week: {week}\n"
        f"date: {lecture_date}\n"
        f"source_pptx: {source_path.name}\n"
        f"source_extracted: {date.today().isoformat()}\n"
        f"slide_count: {slide_count}\n"
        f"related_lab: week-{week:02d}\n"
        "---\n"
    )


def split_existing(text: str) -> tuple[str, str]:
    """Split an existing corpus file into (frontmatter+heading prefix, section_b_text).
    If no Section B exists, returns (text, '')."""
    idx = text.find(SECTION_B_HEADER)
    if idx == -1:
        return text, ""
    return text[:idx], text[idx:]


def write_lecture_file(lecture_n: int, source_path: Path) -> Path:
    slug, title, week, lecture_date = LECTURE_META[lecture_n]
    target = VAULT_LECTURES_DIR / f"{lecture_n:02d}-{slug}.md"
    slide_count, body = extract_pptx(source_path)

    frontmatter = render_frontmatter(
        lecture_n, slug, title, week, lecture_date, source_path, slide_count,
    )
    heading = f"\n# Lecture {lecture_n} — {title}\n\n"
    section_a = render_section_a(lecture_n, source_path, slide_count, body)

    # Preserve existing Section B if file exists.
    if target.exists():
        existing = target.read_text(encoding="utf-8")
        _, section_b = split_existing(existing)
        if not section_b.strip():
            section_b = render_section_b_placeholder()
    else:
        section_b = render_section_b_placeholder()

    out = frontmatter + heading + section_a + "\n---\n\n" + section_b
    target.parent.mkdir(parents=True, exist_ok=True)
    target.write_text(out, encoding="utf-8")
    return target


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--only", type=int, help="Extract only this lecture number")
    args = ap.parse_args()

    if not SOURCE_DIR.exists():
        print(f"ERROR: source dir not found: {SOURCE_DIR}", file=sys.stderr)
        return 1

    targets = [args.only] if args.only else sorted(LECTURE_META.keys())
    for n in targets:
        if n not in LECTURE_META:
            print(f"  SKIP lecture {n} — not in LECTURE_META (out of scope or unknown)", file=sys.stderr)
            continue
        src = find_source_pptx(n)
        if not src:
            print(f"  MISS lecture {n} — no .pptx found in {SOURCE_DIR}", file=sys.stderr)
            continue
        try:
            target = write_lecture_file(n, src)
            print(f"  ok  L{n:02d}  {src.name}  ->  {target.name}")
        except Exception as exc:
            print(f"  FAIL L{n:02d}  {src.name}: {type(exc).__name__}: {exc}", file=sys.stderr)
            return 2
    return 0


if __name__ == "__main__":
    sys.exit(main())
